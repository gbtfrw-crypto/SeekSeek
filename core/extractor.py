"""문서 본문 추출 모듈

다양한 파일 형식에서 텍스트를 추출하여 문자열로 반환한다.
지원 형식: PDF, DOCX, XLSX, PPTX, HWPX, HWP, 일반 텍스트(40여 종)

■ 파일 형식별 파싱 전략
  PDF  : PyMuPDF(fitz) — 페이지 단위 텍스트 추출 (레이아웃 유지)
  DOCX : python-docx  — 문단(paragraph) 단위 추출
  XLSX : openpyxl read_only 모드 — 셀 단위 텍스트 수집
  PPTX : python-pptx → 실패 시 ZIP 직접 파싱 (regex로 <a:t> 태그 추출)
  HWPX : ZIP+XML — Contents/*.xml 내 <t> 태그 추출 (HWP 다음 세대 포맷)
  HWP  : OLE2 바이너리 — olefile + 태그 레코드 파싱 (HWPTAG_PARA_TEXT) (HWP5 포맷)
  텍스트: UTF-8 → CP949 → EUC-KR → Latin-1 순으로 디코딩 시도

■ 선택적 의존성
  외부 라이브러리(fitz, docx, openpyxl, pptx, olefile)는 선택적 의존성으로,
  설치되지 않은 경우 해당 형식만 건너뛰고 None 을 반환한다.
  ImportError는 logger.debug 수준으로 기록하고 상위로 전파하지 않는다.
"""
import os
import re
import zlib
import zipfile
import logging
import xml.etree.ElementTree as ET

import config


logger = logging.getLogger(__name__)

# python-pptx 선택적 임포트. 미설치 시 Presentation=None → ZIP 폴백 경로 사용.
try:
    from pptx import Presentation
except ImportError as e:
    logger.debug("python-pptx ImportError: %s", e)
    Presentation = None

# ── HWP 바이너리 파싱 상수 ────────────────────────────────────────────────────
#
# HWP5 문서 구조 — OLE2(Compound File Binary Format) 컨테이너
# ├─ FileHeader       : 문서 식별·버전·압축 플래그 등 메타데이터
# ├─ DocInfo          : 스타일·폰트·섹션 정의 등 문서 정보
# ├─ BodyText/Section0: 첫 번째 섹션의 본문 태그 레코드 스트림
# ├─ BodyText/Section1: 두 번째 섹션 ...
# └─ ...
#
# 태그 레코드 헤더 (4바이트 little-endian):
#   [tag_id: 10비트] [level: 10비트] [size: 12비트]
#   - tag_id & 0x3FF → 태그 종류  (67 = HWPTAG_PARA_TEXT)
#   - size == 0xFFF  → 확장 크기 (다음 4바이트에 실제 크기)
_HWP_TAG_MASK    = 0x3FF        # tag_id 추출 마스크 (하위 10비트)
_HWP_SIZE_MASK   = 0xFFF        # size 추출 마스크 (상위 12비트)
_HWP_SIZE_SHIFT  = 20           # size 필드 시작 비트 위치
_HWPTAG_PARA_TEXT = 67          # 문단 텍스트 태그 ID (HWPTAG_PARA_TEXT)
_HWP_EXT_SIZE_SENTINEL = 0xFFF  # size가 이 값이면 다음 4바이트가 실제 크기

# HWPTAG_PARA_TEXT 내부 인라인 제어 코드
# 이 코드들은 조판 마크업(change tracking), 내장 그림 등 비텍스트 요소를 나타내며,
# 각 제어 코드 뒤에 12바이트(제어 문자 보조 데이터)가 추가로 따라온다.
_HWP_INLINE_CTRL_CODES = frozenset({1, 2, 3, 11, 12, 14, 15, 16, 17, 18, 21, 22, 23})
_HWP_INLINE_CTRL_EXTRA_BYTES = 12  # 제어 코드 뒤에 오는 고정 12바이트 보조 데이터

# PPTX XML <a:t> 텍스트 추출용 정규식
# PPTX는 ZIP 안에 XML 슬라이드 파일이 있고, 텍스트는 <a:t> 태그 안에 있음
_PPTX_T_RE = re.compile(r"<a:t[^>]*>(.*?)</a:t>", re.DOTALL)
# PPTX 슬라이드 파일명 패턴 (ppt/slides/slide1.xml, slide2.xml, ...)
_PPTX_SLIDE_RE = re.compile(r"ppt/slides/slide\d+\.xml$")


# ── 공개 API ─────────────────────────────────────────────────────────────────

def extract_text(filepath: str) -> str | None:
    """파일에서 텍스트를 추출한다.

    확장자에 따라 적합한 파서로 라우팅한다.
    파싱 오류나 라이브러리 미설치 시 None 을 반환하며, 예외를 외부로 전파하지 않는다.

    Args:
        filepath: 텍스트를 추출할 파일의 절대 경로.

    Returns:
        추출된 텍스트 문자열, 또는 추출 불가 시 None.
    """
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == ".pdf":   return _extract_pdf(filepath)
        if ext == ".docx":  return _extract_docx(filepath)
        if ext == ".xlsx":  return _extract_xlsx(filepath)
        if ext == ".pptx":  return _extract_pptx(filepath)
        if ext == ".hwpx":  return _extract_hwpx(filepath)
        if ext == ".hwp":   return _extract_hwp(filepath)
        return _extract_plain(filepath)
    except Exception as e:
        logger.debug("내용 추출 실패 %s: %s", filepath, e)
        return None


# ── 형식별 파서 ───────────────────────────────────────────────────────────────

def _extract_plain(filepath: str) -> str | None:
    """일반 텍스트 파일을 읽어 문자열로 반환한다.

    인코딩 자동 감지: UTF-8 → CP949 → EUC-KR → Latin-1 순으로 시도한다.
    파일을 바이트로 한 번만 읽고 디코딩만 반복하므로 I/O가 1회로 최소화된다.
    Latin-1은 모든 바이트를 유효한 문자로 매핑하므로 최종 폴백으로 항상 성공한다.
    """
    try:
        with open(filepath, "rb") as f:
            raw = f.read(config.MAX_CONTENT_SIZE)
    except OSError:
        return None
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return None


def _extract_pdf(filepath: str) -> str | None:
    """PyMuPDF(fitz)로 PDF에서 텍스트를 추출한다.

    페이지 단위로 get_text()를 호출하며, 빈 페이지는 결과에서 제외한다.
    fitz가 미설치 상태면 None을 반환하고 경고 로그를 남긴다.
    """
    try:
        import fitz
    except ImportError:
        logger.warning("PyMuPDF 미설치 — PDF 내용 인덱싱 건너뜀")
        return None
    doc = fitz.open(filepath)
    try:
        # 빈 페이지(텍스트 없음)는 제외하고 페이지 텍스트를 리스트로 수집
        texts = [t for page in doc if (t := page.get_text())]
    finally:
        doc.close()
    return "\n".join(texts) if texts else None


def _extract_docx(filepath: str) -> str | None:
    """python-docx로 DOCX에서 문단 텍스트를 추출한다.

    Document.paragraphs를 순회하며 공백이 아닌 문단만 수집한다.
    표(table) 내 셀 텍스트는 현재 추출하지 않는다.
    """
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx 미설치 — DOCX 내용 인덱싱 건너뜀")
        return None
    doc   = Document(filepath)
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(texts) if texts else None


def _extract_xlsx(filepath: str) -> str | None:
    """openpyxl read_only 모드로 XLSX 셀 텍스트를 추출한다.

    read_only=True, data_only=True 설정으로 수식 대신 계산된 값을 읽고
    대용량 파일의 메모리 소비를 최소화한다.
    모든 시트의 셀을 순회하며 None이 아닌 값을 공백으로 연결한다.
    """
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl 미설치 — XLSX 내용 인덱싱 건너뜀")
        return None
    wb    = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            # 행 내 셀 값을 문자열로 변환하여 공백으로 연결, 빈 행은 제외
            row_text = " ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                texts.append(row_text)
    wb.close()
    return "\n".join(texts) if texts else None


def _extract_pptx(filepath: str) -> str | None:
    """PPTX에서 슬라이드 텍스트를 추출한다.

    ■ 추출 전략 (2단계 폴백)
      1차: python-pptx의 Presentation API 사용
           - slide.shapes를 순회하여 text_frame이 있는 도형에서 문단 텍스트 수집
           - 가장 정확하고 텍스트 구조를 잘 보존하지만 대용량 파일에서 느릴 수 있음
      2차: _extract_pptx_fallback() — ZIP 직접 파싱
           - python-pptx가 미설치이거나 파싱 오류 시 사용
           - PPTX는 ZIP이므로 슬라이드 XML을 직접 열어 <a:t> 태그 정규식으로 추출
           - python-pptx보다 가볍지만 텍스트 구조 정보가 없음
    """
    try:
        if Presentation is None:
            # python-pptx가 없으면 바로 ZIP 폴백 경로로 진행
            return _extract_pptx_fallback(filepath)
        prs   = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # text_frame.paragraphs: 도형 내 문단 목록
                    for para in shape.text_frame.paragraphs:
                        if t := para.text.strip():
                            texts.append(t)
        return "\n".join(texts) if texts else None
    except Exception as e:
        # 암호화된 파일, 손상된 ZIP 구조 등 파싱 실패 시 ZIP 폴백 시도
        logger.debug("python-pptx 파싱 실패, ZIP 폴백 시도 %s: %s", filepath, e)
        return _extract_pptx_fallback(filepath)


def _extract_pptx_fallback(filepath: str) -> str | None:
    """PPTX를 ZIP으로 열어 슬라이드 XML에서 <a:t> 텍스트를 직접 추출한다.

    PPTX 파일 구조:
      ppt/slides/slide1.xml, slide2.xml, ... 에 슬라이드별 DrawingML XML이 있음.
      텍스트 내용은 <a:t> 태그 안에 위치한다.

    슬라이드를 정렬된 순서로 처리하여 원본 순서를 유지한다.
    개별 슬라이드 XML 파싱 실패는 해당 슬라이드만 건너뛰고 계속 진행한다.
    """
    texts = []
    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            # 슬라이드 파일만 필터링 후 이름순 정렬(슬라이드 번호 순)
            slide_names = sorted(n for n in zf.namelist() if _PPTX_SLIDE_RE.match(n))
            for name in slide_names:
                try:
                    xml_str = zf.read(name).decode("utf-8", errors="replace")
                    for m in _PPTX_T_RE.finditer(xml_str):
                        if t := m.group(1).strip():
                            texts.append(t)
                except Exception as e:
                    logger.debug("PPTX 슬라이드 XML 읽기 실패 %s/%s: %s", filepath, name, e)
    except Exception as e:
        logger.debug("PPTX ZIP 폴백 실패 %s: %s", filepath, e)
        return None
    return "\n".join(texts) if texts else None


def _extract_hwpx(filepath: str) -> str | None:
    """HWPX(ZIP+XML) 파일에서 <t> 태그 텍스트를 추출한다.

    HWPX 파일 구조:
      HWP 다음 세대 포맷. ZIP 컨테이너 안에 Contents/*.xml 형태로
      섹션별 XML이 저장된다. 텍스트 내용은 네임스페이스 무관 <t> 태그 안에 있음.

    ElementTree의 iter()로 모든 하위 태그를 순회하며, 네임스페이스 접두사를
    제거(split("}")[-1])하여 태그명만 비교한다.
    """
    texts = []
    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            for name in sorted(zf.namelist()):
                # Contents/ 폴더의 .xml 파일만 처리 (섹션 본문 XML)
                if not (name.startswith("Contents/") and name.endswith(".xml")):
                    continue
                try:
                    root = ET.fromstring(zf.read(name))
                    for elem in root.iter():
                        # {namespace}tagname 형식에서 tagname만 추출
                        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                        if tag == "t" and elem.text and elem.text.strip():
                            texts.append(elem.text.strip())
                except (ET.ParseError, UnicodeDecodeError):
                    # 특정 섹션 XML이 손상되었어도 나머지 섹션 계속 처리
                    continue
    except zipfile.BadZipFile:
        # HWPX 확장자이지만 올바른 ZIP 구조가 아닌 파일
        return None
    return "\n".join(texts) if texts else None


def _hwp_decompress(data: bytes, filepath: str, stream_name: str) -> bytes | None:
    """HWP BodyText 스트림의 zlib 압축을 해제한다.

    ■ HWP5 압축 스펙
      FileHeader의 flags 필드(오프셋 +36, 4바이트, little-endian) 중
      비트 0이 1이면 BodyText 스트림이 압축됨.
      HWP5 표준은 headerless raw deflate(wbits=-15)를 사용하지만,
      일부 한글 버전은 zlib 헤더(0x78 9C / 0x78 DA)를 포함하거나
      gzip 헤더를 사용하는 비표준 압축 방식으로 저장한다.

    ■ 3단계 디코딩 전략 (실패 시 다음 단계 시도)
      1차: wbits=-15 → raw deflate (HWP5 표준)
      2차: wbits=+15 → zlib 헤더 포함 wrapped deflate
           "invalid distance too far back" 오류는 raw 파서가 zlib 헤더 2바이트를
           데이터로 잘못 해석해 LZ77 distance 테이블이 틀어질 때 전형적으로 발생
      3차: wbits=47 → gzip/zlib 자동 감지 (gzip 헤더 포함 파일 대응)

    Args:
        data:        BodyText 스트림 원본 바이트.
        filepath:    로그용 파일 경로.
        stream_name: 로그용 OLE 스트림 경로 (예: "BodyText/Section0").

    Returns:
        압축 해제된 바이트, 또는 모든 시도 실패 시 None.
    """
    # 1차 시도: HWP5 표준 — headerless raw deflate
    try:
        return zlib.decompress(data, -15)
    except zlib.error:
        pass

    # 2차 시도: zlib 헤더(0x78 9C / 0x78 DA 등) 포함된 wrapped deflate
    try:
        return zlib.decompress(data, 15)
    except zlib.error:
        pass

    # 3차 시도: gzip/zlib 자동 감지 (wbits=47 = 32+15, MAX_WBITS with auto-detect)
    try:
        return zlib.decompress(data, 47)
    except zlib.error as e:
        logger.debug("HWP decompress 실패 (%s/%s): %s", filepath, stream_name, e)
        return None


def _extract_hwp(filepath: str) -> str | None:
    """레거시 HWP(OLE2 바이너리) 파일에서 텍스트를 추출한다.

    ■ 처리 흐름
      1. olefile로 OLE2 컨테이너를 열고 FileHeader 스트림에서 압축 플래그를 확인
      2. BodyText/Section* 스트림을 순회
      3. compressed=True 이면 _hwp_decompress()로 압축 해제
      4. 해제된(또는 원본) 바이트 스트림에서 HWPTAG_PARA_TEXT 레코드 추출
      5. 레코드 데이터를 UTF-16LE로 디코딩하여 텍스트 수집

    ■ 압축 플래그 위치
      FileHeader 스트림 오프셋 36~39바이트(4바이트, little-endian)의 bit 0
      (0 = 비압축, 1 = 압축)
    """
    try:
        import olefile
    except ImportError:
        logger.warning("olefile 미설치 — HWP 내용 인덱싱 건너뜀")
        return None

    try:
        ole = olefile.OleFileIO(filepath)
    except Exception as e:
        logger.debug("OLE 파일 열기 실패 (%s): %s", filepath, e)
        return None

    # FileHeader에서 압축 플래그(bit 0) 확인
    # HWP5 스펙: FileHeader 크기 최소 256바이트, flags는 +36 오프셋에 있음
    compressed = False
    if ole.exists("FileHeader"):
        header_data = ole.openstream("FileHeader").read()
        if len(header_data) >= 40:
            flags = int.from_bytes(header_data[36:40], "little")
            compressed = bool(flags & 0x1)
    logger.debug("HWP 압축 여부 (%s): %s", filepath, compressed)

    texts = []
    try:
        for stream_name in ole.listdir():
            # BodyText/Section* 스트림만 처리 (본문 데이터)
            if not "/".join(stream_name).startswith("BodyText/Section"):
                continue
            data = ole.openstream(stream_name).read()
            if compressed:
                # 압축 파일: 해제 실패 시 해당 섹션 스킵 (깨진 데이터 파싱 방지)
                data = _hwp_decompress(data, filepath, "/".join(stream_name))
                if data is None:
                    continue
            if text := _extract_text_from_hwp_body(data):
                texts.append(text)
    finally:
        ole.close()
    return "\n".join(texts) if texts else None


# ── HWP 바이너리 파서 ─────────────────────────────────────────────────────────

def _extract_text_from_hwp_body(data: bytes) -> str:
    """HWP BodyText 스트림에서 HWPTAG_PARA_TEXT(태그 ID=67) 레코드를 추출한다.

    ■ 태그 레코드 순회 알고리즘
      1) 현재 오프셋에서 4바이트 헤더 읽기
      2) 헤더에서 tag_id(하위 10비트)와 size(상위 12비트) 추출
      3) size == 0xFFF → 확장 크기 모드: 다음 4바이트가 실제 페이로드 크기
         (헤더 소비 바이트: 일반 4바이트, 확장 8바이트)
      4) tag_id == 67(HWPTAG_PARA_TEXT) → _decode_hwp_para_text()로 텍스트 디코딩
      5) 그 외 태그 → 페이로드만큼 오프셋 전진하고 건너뜀
      6) 스트림 끝까지 반복

    Args:
        data: 압축 해제된 BodyText 스트림 바이트.

    Returns:
        추출된 모든 문단을 줄바꿈으로 연결한 문자열.
    """
    texts = []
    i     = 0
    while i < len(data) - 4:
        header = int.from_bytes(data[i:i + 4], "little")
        tag_id = header & _HWP_TAG_MASK
        size   = (header >> _HWP_SIZE_SHIFT) & _HWP_SIZE_MASK

        if size == _HWP_EXT_SIZE_SENTINEL:
            # 확장 크기: 다음 4바이트에 실제 페이로드 크기가 있음
            if i + 8 > len(data):
                break
            size  = int.from_bytes(data[i + 4:i + 8], "little")
            i    += 8  # 4바이트 기본 헤더 + 4바이트 확장 크기 소비
        else:
            i += 4  # 기본 4바이트 헤더 소비

        if i + size > len(data):
            break  # 페이로드가 버퍼 경계를 넘으면 중단

        if tag_id == _HWPTAG_PARA_TEXT:
            if text := _decode_hwp_para_text(data[i:i + size]).strip():
                texts.append(text)
        i += size  # 페이로드 건너뛰기

    return "\n".join(texts)


def _decode_hwp_para_text(chunk: bytes) -> str:
    """HWP PARA_TEXT 레코드를 UTF-16LE 코드 단위로 디코딩한다.

    ■ PARA_TEXT 인코딩 규칙 (UTF-16LE, 2바이트 단위)
      - code >= 0x0020 : 일반 유니코드 문자 → chr(code)
      - code == 0x000D, 0x000A : 줄바꿈 (CR/LF)
      - code == 0x0009 : 탭 → 공백으로 대체
      - code in _HWP_INLINE_CTRL_CODES : 인라인 제어 코드
        → 뒤따르는 12바이트(보조 데이터)를 함께 건너뜀

    Args:
        chunk: HWPTAG_PARA_TEXT 레코드의 페이로드 바이트.

    Returns:
        디코딩된 문자열.
    """
    chars = []
    j     = 0
    while j < len(chunk) - 1:
        code = int.from_bytes(chunk[j:j + 2], "little")
        j   += 2
        if code >= 0x0020:
            chars.append(chr(code))
        elif code in (0x0D, 0x0A):
            chars.append("\n")
        elif code == 0x09:
            chars.append(" ")
        elif code in _HWP_INLINE_CTRL_CODES:
            # 제어 코드 뒤에 오는 12바이트 보조 데이터 건너뜀
            j += _HWP_INLINE_CTRL_EXTRA_BYTES
    return "".join(chars)
